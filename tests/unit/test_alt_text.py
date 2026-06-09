"""Unit tests for alt_text rule."""

import io

from PIL import Image  # type: ignore[import-untyped]
from pptx import Presentation  # type: ignore[import-untyped]
from pptx.chart.data import ChartData  # type: ignore[import-untyped]
from pptx.enum.chart import XL_CHART_TYPE  # type: ignore[import-untyped]
from pptx.oxml import parse_xml  # type: ignore[import-untyped]
from pptx.oxml.ns import nsdecls  # type: ignore[import-untyped]
from pptx.util import Inches  # type: ignore[import-untyped]

from a11yfix.ooxml.image_extract import extract_image_for_finding
from a11yfix.ooxml.pptx_reader import open_pptx
from a11yfix.rules.alt_text import AltTextQualityRule, AltTextRule


def test_alt_missing_detected(pptx_no_alt):
    doc = open_pptx(pptx_no_alt)
    findings = list(AltTextRule().detect(doc))
    assert any(f.rule_id == "alt-text-missing" for f in findings)


def test_alt_present_not_flagged(pptx_with_alt):
    doc = open_pptx(pptx_with_alt)
    findings = list(AltTextRule().detect(doc))
    # The fixture has at least one image with alt; we should not flag it
    assert not findings
    assert not list(AltTextQualityRule().detect(doc))


def test_pptx_text_shapes_are_not_image_alt_findings(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])  # Title Only
    slide.shapes.title.text = "A real slide title"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    textbox.text = "Ordinary body text in a text box"
    path = tmp_path / "text_shapes_only.pptx"
    pres.save(path)

    doc = open_pptx(path)
    findings = list(AltTextRule().detect(doc))

    assert findings == []


def test_pptx_real_picture_without_alt_is_flagged_amid_text_shapes(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[5])
    slide.shapes.title.text = "A real slide title"
    textbox = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    textbox.text = "Ordinary body text in a text box"

    buf = io.BytesIO()
    Image.new("RGB", (50, 50), color="blue").save(buf, format="PNG")
    buf.seek(0)
    slide.shapes.add_picture(buf, Inches(1), Inches(3), Inches(2), Inches(2))

    path = tmp_path / "picture_and_text.pptx"
    pres.save(path)

    doc = open_pptx(path)
    findings = list(AltTextRule().detect(doc))

    assert len(findings) == 1
    assert findings[0].officecli_path.startswith("/slide[1]/picture[@id=")


def test_pptx_image_filled_shape_without_alt_is_flagged_and_extractable(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])

    buf = io.BytesIO()
    Image.new("RGB", (50, 50), color="green").save(buf, format="PNG")
    buf.seek(0)
    source_pic = slide.shapes.add_picture(buf, Inches(8), Inches(6), Inches(1), Inches(1))
    source_pic._element.nvPicPr.cNvPr.set("descr", "source image")

    embed = ""
    for el in source_pic._element.iter():
        if el.tag.endswith("}blip"):
            embed = el.get(
                "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
            )
            break
    assert embed

    image_shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(3), Inches(2))
    blip_fill = parse_xml(
        f'<a:blipFill {nsdecls("a", "r")}><a:blip r:embed="{embed}"/></a:blipFill>'
    )
    image_shape._element.spPr.append(blip_fill)

    path = tmp_path / "image_filled_shape.pptx"
    pres.save(path)

    doc = open_pptx(path)
    findings = list(AltTextRule().detect(doc))

    assert len(findings) == 1
    assert findings[0].officecli_path.startswith("/slide[1]/shape[@id=")
    extracted = extract_image_for_finding(doc, findings[0])
    assert extracted is not None
    assert extracted[1] == "image/png"


def test_pptx_auto_generated_alt_is_quality_warning_not_missing(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    buf = io.BytesIO()
    Image.new("RGB", (50, 50), color="purple").save(buf, format="PNG")
    buf.seek(0)
    pic = slide.shapes.add_picture(buf, Inches(1), Inches(1), Inches(2), Inches(2))
    pic._element.nvPicPr.cNvPr.set("descr", "Chart\n\nDescription automatically generated")

    path = tmp_path / "auto_generated_alt.pptx"
    pres.save(path)

    doc = open_pptx(path)
    missing = list(AltTextRule().detect(doc))
    quality = list(AltTextQualityRule().detect(doc))

    assert missing == []
    assert len(quality) == 1
    assert quality[0].rule_id == "alt-text-generic"
    assert quality[0].severity.value == "warning"
    assert quality[0].current_value == "Chart\n\nDescription automatically generated"
    assert quality[0].extra["reason"] == "office_auto_generated"
    assert AltTextQualityRule().fix_single_shot(quality[0], doc).kind == "alt-text"


def test_pptx_local_path_alt_is_quality_warning(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    buf = io.BytesIO()
    Image.new("RGB", (50, 50), color="orange").save(buf, format="PNG")
    buf.seek(0)
    pic = slide.shapes.add_picture(buf, Inches(1), Inches(1), Inches(2), Inches(2))
    pic._element.nvPicPr.cNvPr.set("descr", r"C:\clients\aps\AP&S Stacked.png")

    path = tmp_path / "path_alt.pptx"
    pres.save(path)

    doc = open_pptx(path)
    missing = list(AltTextRule().detect(doc))
    quality = list(AltTextQualityRule().detect(doc))

    assert missing == []
    assert len(quality) == 1
    assert quality[0].extra["reason"] == "local_file_path"


def test_pptx_picture_auto_name_stays_missing(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    buf = io.BytesIO()
    Image.new("RGB", (50, 50), color="yellow").save(buf, format="PNG")
    buf.seek(0)
    pic = slide.shapes.add_picture(buf, Inches(1), Inches(1), Inches(2), Inches(2))
    pic._element.nvPicPr.cNvPr.set("descr", "Picture 4")

    path = tmp_path / "picture_auto_name.pptx"
    pres.save(path)

    doc = open_pptx(path)
    missing = list(AltTextRule().detect(doc))
    quality = list(AltTextQualityRule().detect(doc))

    assert len(missing) == 1
    assert missing[0].rule_id == "alt-text-missing"
    assert quality == []


def test_pptx_chart_without_alt_is_flagged_as_object_alt(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    data = ChartData()
    data.categories = ["A", "B"]
    data.add_series("Series", (1, 2))
    slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        data,
    )
    path = tmp_path / "chart_no_alt.pptx"
    pres.save(path)

    findings = list(AltTextRule().detect(open_pptx(path)))

    assert len(findings) == 1
    assert findings[0].officecli_path.startswith("/slide[1]/chart[@id=")
    assert findings[0].extra["shape_kind"] == "chart"


def test_pptx_chart_with_alt_is_not_flagged(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    data = ChartData()
    data.categories = ["A", "B"]
    data.add_series("Series", (1, 2))
    chart_shape = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(1),
        Inches(1),
        Inches(4),
        Inches(3),
        data,
    )
    chart_shape._element.nvGraphicFramePr.cNvPr.set("descr", "Bar chart comparing A and B")
    path = tmp_path / "chart_with_alt.pptx"
    pres.save(path)

    findings = list(AltTextRule().detect(open_pptx(path)))

    assert findings == []


def test_pptx_grouped_picture_path_includes_group_scope(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    group = parse_xml(
        f"""
        <p:grpSp {nsdecls("p", "a", "r")}>
          <p:nvGrpSpPr>
            <p:cNvPr id="99" name="Grouped objects"/>
            <p:cNvGrpSpPr/>
            <p:nvPr/>
          </p:nvGrpSpPr>
          <p:grpSpPr>
            <a:xfrm>
              <a:off x="0" y="0"/><a:ext cx="1" cy="1"/>
              <a:chOff x="0" y="0"/><a:chExt cx="1" cy="1"/>
            </a:xfrm>
          </p:grpSpPr>
          <p:pic>
            <p:nvPicPr>
              <p:cNvPr id="100" name="Grouped picture"/>
              <p:cNvPicPr/>
              <p:nvPr/>
            </p:nvPicPr>
            <p:blipFill/>
            <p:spPr/>
          </p:pic>
        </p:grpSp>
        """
    )
    slide._element.spTree.append(group)
    path = tmp_path / "grouped_picture.pptx"
    pres.save(path)

    findings = list(AltTextRule().detect(open_pptx(path)))
    picture = next(f for f in findings if f.extra["shape_kind"] == "picture")

    assert picture.officecli_path == "/slide[1]/group[@id=99]/picture[@id=100]"


def test_pptx_nested_grouped_picture_path_includes_full_group_scope(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    group = parse_xml(
        f"""
        <p:grpSp {nsdecls("p", "a", "r")}>
          <p:nvGrpSpPr>
            <p:cNvPr id="90" name="Outer group"/>
            <p:cNvGrpSpPr/>
            <p:nvPr/>
          </p:nvGrpSpPr>
          <p:grpSpPr>
            <a:xfrm>
              <a:off x="0" y="0"/><a:ext cx="1" cy="1"/>
              <a:chOff x="0" y="0"/><a:chExt cx="1" cy="1"/>
            </a:xfrm>
          </p:grpSpPr>
          <p:grpSp>
            <p:nvGrpSpPr>
              <p:cNvPr id="99" name="Inner group"/>
              <p:cNvGrpSpPr/>
              <p:nvPr/>
            </p:nvGrpSpPr>
            <p:grpSpPr>
              <a:xfrm>
                <a:off x="0" y="0"/><a:ext cx="1" cy="1"/>
                <a:chOff x="0" y="0"/><a:chExt cx="1" cy="1"/>
              </a:xfrm>
            </p:grpSpPr>
            <p:pic>
              <p:nvPicPr>
                <p:cNvPr id="100" name="Nested grouped picture"/>
                <p:cNvPicPr/>
                <p:nvPr/>
              </p:nvPicPr>
              <p:blipFill/>
              <p:spPr/>
            </p:pic>
          </p:grpSp>
        </p:grpSp>
        """
    )
    slide._element.spTree.append(group)
    path = tmp_path / "nested_grouped_picture.pptx"
    pres.save(path)

    findings = list(AltTextRule().detect(open_pptx(path)))
    picture = next(f for f in findings if f.extra["shape_kind"] == "picture")

    assert picture.officecli_path == (
        "/slide[1]/group[@id=90]/group[@id=99]/picture[@id=100]"
    )


def test_pptx_picture_inside_idless_group_is_not_flattened_to_slide_path(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    group = parse_xml(
        f"""
        <p:grpSp {nsdecls("p", "a", "r")}>
          <p:nvGrpSpPr>
            <p:cNvPr name="Idless group"/>
            <p:cNvGrpSpPr/>
            <p:nvPr/>
          </p:nvGrpSpPr>
          <p:grpSpPr>
            <a:xfrm>
              <a:off x="0" y="0"/><a:ext cx="1" cy="1"/>
              <a:chOff x="0" y="0"/><a:chExt cx="1" cy="1"/>
            </a:xfrm>
          </p:grpSpPr>
          <p:pic>
            <p:nvPicPr>
              <p:cNvPr id="100" name="Grouped picture"/>
              <p:cNvPicPr/>
              <p:nvPr/>
            </p:nvPicPr>
            <p:blipFill/>
            <p:spPr/>
          </p:pic>
        </p:grpSp>
        """
    )
    slide._element.spTree.append(group)
    path = tmp_path / "idless_grouped_picture.pptx"
    pres.save(path)

    findings = list(AltTextRule().detect(open_pptx(path)))

    assert [f.officecli_path for f in findings] == []


def test_pptx_smartart_is_scan_only_not_fixable(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    smartart = parse_xml(
        f"""
        <p:graphicFrame {nsdecls("p", "a", "r")}>
          <p:nvGraphicFramePr>
            <p:cNvPr id="77" name="SmartArt"/>
            <p:cNvGraphicFramePr/>
            <p:nvPr/>
          </p:nvGraphicFramePr>
          <p:xfrm>
            <a:off x="0" y="0"/><a:ext cx="1" cy="1"/>
          </p:xfrm>
          <a:graphic>
            <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/diagram"/>
          </a:graphic>
        </p:graphicFrame>
        """
    )
    slide._element.spTree.append(smartart)
    path = tmp_path / "smartart.pptx"
    pres.save(path)

    doc = open_pptx(path)
    finding = next(iter(AltTextRule().detect(doc)))

    assert finding.officecli_path == "/slide[1]"
    assert finding.extra["shape_kind"] == "smartArt"
    assert finding.why_human_needed is not None
    assert AltTextRule().fix_single_shot(finding, doc) is None
