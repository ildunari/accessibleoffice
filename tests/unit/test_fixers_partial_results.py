"""Regression tests for partial officecli batch result handling."""

from __future__ import annotations

from pathlib import Path

from docx import Document  # type: ignore[import-untyped]
from PIL import Image  # type: ignore[import-untyped]
from pptx import Presentation  # type: ignore[import-untyped]
from pptx.util import Inches  # type: ignore[import-untyped]

from a11yfix.fixers import deterministic, single_shot
from a11yfix.manifest import FileFormat, Finding, Severity
from a11yfix.ooxml.docx_reader import open_docx
from a11yfix.ooxml.officecli import BatchResult, ValidationResult
from a11yfix.ooxml.pptx_reader import open_pptx
from a11yfix.rules.alt_text import AltTextRule
from a11yfix.rules.base import OfficecliOp, RuleMeta, SingleShotFix


class FakeDoc:
    file_format = FileFormat.PPTX

    def __init__(self, path: Path) -> None:
        self.path = str(path)

    def root_xml(self):
        return None


class FakeRule:
    meta = RuleMeta(
        rule_id="fake-rule",
        severity=Severity.ERROR,
        formats={FileFormat.PPTX},
        wcag_sc=[],
        plain_impact="fake",
    )

    def fix_deterministic(self, finding, doc):
        return [OfficecliOp(verb="set", path=finding.officecli_path, props={"x": "y"})]

    def fix_single_shot(self, finding, doc):
        return SingleShotFix(kind="link-text", finding=finding)


class ShortResultClient:
    backup_path = None

    def __init__(self, path, **kwargs):
        self.path = path

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, tb):
        return None

    def batch(self, ops):
        return BatchResult(success=False, per_op=[{"ok": True}])

    def validate(self):
        return ValidationResult(status="ok")


class FakeAdapter:
    name = "fake"
    calls = 0

    def suggest_link_text(self, url, surrounding_text):
        self.calls += 1

        class Result:
            text = "descriptive link"
            confidence = 0.9
            model = "fake"

        return Result()


class LowConfidenceAdapter:
    name = "low-confidence"

    def suggest_link_text(self, url, surrounding_text):
        class Result:
            text = "weak link text"
            confidence = 0.2
            model = "fake"

        return Result()


class UnclearAdapter:
    name = "unclear"

    def suggest_link_text(self, url, surrounding_text):
        class Result:
            text = "UNCLEAR"
            confidence = 0.9
            model = "fake"

        return Result()


def _finding(i: int) -> Finding:
    return Finding(
        id=f"f{i}",
        rule_id="fake-rule",
        severity=Severity.ERROR,
        wcag_sc=[],
        officecli_path=f"/slide[1]/picture[@id={i}]",
        extra={"url": "https://example.com", "shape_text": "click"},
    )


def test_deterministic_short_officecli_results_defer_missing_ops(tmp_path, monkeypatch):
    doc = FakeDoc(tmp_path / "deck.pptx")
    findings = [_finding(1), _finding(2)]

    monkeypatch.setitem(deterministic.REGISTRY, "fake-rule", FakeRule())
    monkeypatch.setattr(deterministic, "OfficecliClient", ShortResultClient)

    result = deterministic.apply_deterministic_fixes(findings, doc)

    assert [f.finding_id for f in result.applied] == ["f1"]
    assert [f.id for f in result.deferred] == ["f2"]


class NoopRestoreClient:
    """Validation passes but every op fails — the open/save round-trip still
    mutated the file, so the fixer must restore the pristine backup."""

    backup_path = "/tmp/whatever.bak"

    def __init__(self, path, **kwargs):
        self.path = path
        self.restored = False

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc, tb):
        return None

    def batch(self, ops):
        return BatchResult(success=False, per_op=[{"ok": False} for _ in ops])

    def validate(self):
        return ValidationResult(status="ok")

    def restore_from_backup(self):
        self.restored = True


def test_deterministic_restores_backup_when_no_fix_applied(tmp_path, monkeypatch):
    doc = FakeDoc(tmp_path / "deck.pptx")
    findings = [_finding(1), _finding(2)]
    clients: list[NoopRestoreClient] = []

    def make_client(path, **kwargs):
        c = NoopRestoreClient(path, **kwargs)
        clients.append(c)
        return c

    monkeypatch.setitem(deterministic.REGISTRY, "fake-rule", FakeRule())
    monkeypatch.setattr(deterministic, "OfficecliClient", make_client)

    result = deterministic.apply_deterministic_fixes(findings, doc)

    assert result.applied == []
    assert [f.id for f in result.deferred] == ["f1", "f2"]
    assert clients and clients[0].restored, "no-op run must restore the backup byte-for-byte"


def test_single_shot_short_officecli_results_defer_missing_ops(tmp_path, monkeypatch):
    doc = FakeDoc(tmp_path / "deck.pptx")
    findings = [_finding(1), _finding(2)]

    monkeypatch.setitem(single_shot.REGISTRY, "fake-rule", FakeRule())
    monkeypatch.setattr(single_shot, "OfficecliClient", ShortResultClient)

    result = single_shot.apply_single_shot_fixes(findings, doc, FakeAdapter())

    assert [f.finding_id for f in result.applied] == ["f1"]
    assert [f.id for f in result.deferred] == ["f2"]


def test_single_shot_zero_cost_cap_defers_without_adapter_call(tmp_path, monkeypatch):
    doc = FakeDoc(tmp_path / "deck.pptx")
    findings = [_finding(1)]
    adapter = FakeAdapter()

    monkeypatch.setitem(single_shot.REGISTRY, "fake-rule", FakeRule())

    result = single_shot.apply_single_shot_fixes(
        findings,
        doc,
        adapter,
        max_cost_total_usd=0,
    )

    assert adapter.calls == 0
    assert result.applied == []
    assert [f.id for f in result.deferred] == ["f1"]
    assert result.deferred[0].why_human_needed == "stage-3 deferred: batch cost cap reached"


def test_single_shot_low_confidence_deferral_has_reason(tmp_path, monkeypatch):
    doc = FakeDoc(tmp_path / "deck.pptx")
    findings = [_finding(1)]

    monkeypatch.setitem(single_shot.REGISTRY, "fake-rule", FakeRule())
    monkeypatch.setattr(single_shot, "_cache_get", lambda payload: None)
    monkeypatch.setattr(single_shot, "_cache_put", lambda payload, value: None)

    result = single_shot.apply_single_shot_fixes(findings, doc, LowConfidenceAdapter())

    assert result.applied == []
    assert result.deferred[0].why_human_needed == "stage-3 deferred: low confidence (0.20 < 0.60)"


def test_single_shot_unclear_deferral_has_reason(tmp_path, monkeypatch):
    doc = FakeDoc(tmp_path / "deck.pptx")
    findings = [_finding(1)]

    monkeypatch.setitem(single_shot.REGISTRY, "fake-rule", FakeRule())
    monkeypatch.setattr(single_shot, "_cache_get", lambda payload: None)
    monkeypatch.setattr(single_shot, "_cache_put", lambda payload, value: None)

    result = single_shot.apply_single_shot_fixes(findings, doc, UnclearAdapter())

    assert result.applied == []
    assert result.deferred[0].why_human_needed == "stage-3 deferred: model returned UNCLEAR"


def test_single_shot_pptx_alt_text_emits_officecli_alt_op(tmp_path, monkeypatch):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    img_path = tmp_path / "image.png"
    Image.new("RGB", (50, 50), color="blue").save(img_path)
    slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(2), Inches(2))
    deck = tmp_path / "deck.pptx"
    pres.save(deck)

    doc = open_pptx(deck)
    finding = next(iter(AltTextRule().detect(doc)))
    captured: list[OfficecliOp] = []

    class AltAdapter:
        name = "fake-alt"

        def describe_image(self, image_bytes, *, max_chars, context):
            class Result:
                text = "Blue square image"
                confidence = 0.95
                model = "fake-alt"

            return Result()

    class CapturingClient:
        backup_path = None

        def __init__(self, path, **kwargs):
            self.path = path

        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, tb):
            return None

        def batch(self, ops):
            captured.extend(ops)
            return BatchResult(success=True, per_op=[{"ok": True}])

        def validate(self):
            return ValidationResult(status="ok")

    monkeypatch.setattr(single_shot, "OfficecliClient", CapturingClient)
    monkeypatch.setattr(single_shot, "_cache_get", lambda payload: None)
    monkeypatch.setattr(single_shot, "_cache_put", lambda payload, value: None)

    result = single_shot.apply_single_shot_fixes([finding], doc, AltAdapter())

    assert [fix.finding_id for fix in result.applied] == [finding.id]
    assert result.deferred == []
    assert len(captured) == 1
    assert captured[0].verb == "set"
    assert captured[0].path == finding.officecli_path
    assert captured[0].props == {"alt": "Blue square image"}


def test_single_shot_docx_alt_text_emits_officecli_alt_op(tmp_path, monkeypatch):
    img_path = tmp_path / "image.png"
    Image.new("RGB", (50, 50), color="green").save(img_path)
    docx = Document()
    docx.add_picture(str(img_path))
    path = tmp_path / "doc.docx"
    docx.save(path)

    doc = open_docx(path)
    finding = next(iter(AltTextRule().detect(doc)))
    captured: list[OfficecliOp] = []

    class AltAdapter:
        name = "fake-alt"

        def describe_image(self, image_bytes, *, max_chars, context):
            class Result:
                text = "Green square image"
                confidence = 0.95
                model = "fake-alt"

            return Result()

    class CapturingClient:
        backup_path = None

        def __init__(self, path, **kwargs):
            self.path = path

        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, tb):
            return None

        def batch(self, ops):
            captured.extend(ops)
            return BatchResult(success=True, per_op=[{"ok": True}])

        def validate(self):
            return ValidationResult(status="ok")

    monkeypatch.setattr(single_shot, "OfficecliClient", CapturingClient)
    monkeypatch.setattr(single_shot, "_cache_get", lambda payload: None)
    monkeypatch.setattr(single_shot, "_cache_put", lambda payload, value: None)

    result = single_shot.apply_single_shot_fixes([finding], doc, AltAdapter())

    assert [fix.finding_id for fix in result.applied] == [finding.id]
    assert result.deferred == []
    assert len(captured) == 1
    assert captured[0].verb == "set"
    assert captured[0].path == finding.officecli_path
    assert captured[0].props == {"alt": "Green square image"}


def test_single_shot_restores_backup_when_no_fix_applied(tmp_path, monkeypatch):
    """Stage 3 mirrors stage 2's no-op guard: if officecli applied none of the
    ops, the open/save round-trip still rewrote the package, so the post-
    stage-2 backup must be restored to keep the bytes untouched."""
    doc = FakeDoc(tmp_path / "deck.pptx")
    findings = [_finding(1), _finding(2)]
    clients: list[NoopRestoreClient] = []

    def make_client(path, **kwargs):
        c = NoopRestoreClient(path, **kwargs)
        clients.append(c)
        return c

    monkeypatch.setattr(single_shot, "CACHE_DIR", tmp_path / "cache")
    monkeypatch.setitem(single_shot.REGISTRY, "fake-rule", FakeRule())
    monkeypatch.setattr(single_shot, "OfficecliClient", make_client)

    result = single_shot.apply_single_shot_fixes(findings, doc, FakeAdapter())

    assert result.applied == []
    assert {f.id for f in result.deferred} == {"f1", "f2"}
    assert clients and clients[0].restored, "no-op stage 3 must restore the backup"


def test_slide_title_cache_key_covers_full_slide_text(tmp_path, monkeypatch):
    """Two slides sharing a layout and the same first 200 chars of text must
    NOT share a cached title — the key must cover the full text the model sees
    (the cache is persistent, so a truncated key leaks titles across slides
    and even across decks)."""
    doc = FakeDoc(tmp_path / "deck.pptx")
    boilerplate = "shared boilerplate opening " * 10  # > 200 chars
    slide_texts = {1: boilerplate + "ends about mitosis", 2: boilerplate + "ends about meiosis"}

    monkeypatch.setattr(single_shot, "CACHE_DIR", tmp_path / "cache")
    monkeypatch.setattr(
        single_shot, "_slide_text", lambda doc, idx: (slide_texts[idx], "Title Layout")
    )

    calls: list[str] = []

    class TitleAdapter:
        name = "fake-title"

        def suggest_slide_title(self, slide_text, layout):
            calls.append(slide_text)

            class Result:
                text = f"Generated title {len(calls)}"
                confidence = 0.9
                model = "fake-title"

            return Result()

    class TitleRule(FakeRule):
        def fix_single_shot(self, finding, doc):
            return SingleShotFix(kind="slide-title", finding=finding)

    class OkClient:
        backup_path = None

        def __init__(self, path, **kwargs):
            self.path = path

        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, tb):
            return None

        def batch(self, ops):
            return BatchResult(success=True, per_op=[{"ok": True} for _ in ops])

        def validate(self):
            return ValidationResult(status="ok")

    monkeypatch.setitem(single_shot.REGISTRY, "fake-rule", TitleRule())
    monkeypatch.setattr(single_shot, "OfficecliClient", OkClient)

    f1, f2 = _finding(1), _finding(2)
    f1.extra["slide_index"] = 1
    f2.extra["slide_index"] = 2

    result = single_shot.apply_single_shot_fixes([f1, f2], doc, TitleAdapter())

    assert len(calls) == 2, "second slide must not reuse the first slide's cached title"
    assert {fx.after for fx in result.applied} == {"Generated title 1", "Generated title 2"}
