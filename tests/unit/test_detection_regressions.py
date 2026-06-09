"""Regression tests for detection-rule fixes (false positives/negatives)."""

from pptx.oxml import parse_xml  # type: ignore[import-untyped]
from pptx.oxml.ns import nsdecls  # type: ignore[import-untyped]

from a11yfix.ooxml.docx_reader import open_docx
from a11yfix.ooxml.pptx_reader import open_pptx
from a11yfix.rules.decorative_flag import DecorativeFlagRule
from a11yfix.rules.heading_structure import _looks_like_fake_heading
from a11yfix.rules.list_semantics import ListSemanticsRule
from a11yfix.rules.reading_order import _shape_position
from a11yfix.rules.table_headers import TableHeaderRule


def _ppt_table(tmp_path, first_row_value: str):
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    path = tmp_path / f"table_{first_row_value}.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1))
    shape.table.cell(0, 0).text = "A"
    shape.table.cell(0, 1).text = "B"
    shape._element.graphic.graphicData.tbl.tblPr.set("firstRow", first_row_value)
    pres.save(path)
    return path


def test_first_row_true_lexical_form_not_flagged(tmp_path):
    # firstRow is xsd:boolean: "true" is just as valid as "1".
    path = _ppt_table(tmp_path, "true")
    findings = list(TableHeaderRule().detect(open_pptx(path)))
    assert findings == []


def test_first_row_one_not_flagged(tmp_path):
    path = _ppt_table(tmp_path, "1")
    findings = list(TableHeaderRule().detect(open_pptx(path)))
    assert findings == []


def test_first_row_false_flagged(tmp_path):
    path = _ppt_table(tmp_path, "false")
    findings = list(TableHeaderRule().detect(open_pptx(path)))
    assert len(findings) == 1


# ---- heading_structure: <w:b w:val="0"/> is bold-OFF -------------------------


def _docx_paragraph_with_bold_off():
    from docx import Document  # type: ignore[import-untyped]

    docx = Document()
    p = docx.add_paragraph()
    run = p.add_run("Short bold-off text")
    run.font.bold = False  # writes <w:b w:val="0"/>
    return p._p


def test_explicit_bold_off_not_a_fake_heading():
    p = _docx_paragraph_with_bold_off()
    assert _looks_like_fake_heading(p) is False


# ---- color_contrast: malformed numeric attributes don't abort detection -----


def test_malformed_sz_and_lummod_do_not_abort_detection(tmp_path):
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    from a11yfix.rules.color_contrast import ColorContrastRule

    path = tmp_path / "malformed.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    box = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run_el = parse_xml(
        f"""
        <a:r {nsdecls("a")}>
          <a:rPr lang="en-US" sz="garbage">
            <a:solidFill>
              <a:srgbClr val="777777">
                <a:lumMod val="not-a-number"/>
              </a:srgbClr>
            </a:solidFill>
          </a:rPr>
          <a:t>low contrast maybe</a:t>
        </a:r>
        """
    )
    box.text_frame.paragraphs[0]._p.append(run_el)
    # Solid white shape background so contrast is computable.
    spPr = box._element.spPr
    spPr.append(
        parse_xml(
            f"""
            <a:solidFill {nsdecls("a")}><a:srgbClr val="FFFFFF"/></a:solidFill>
            """
        )
    )
    pres.save(path)

    # Must not raise; the malformed values fall back to defaults.
    findings = list(ColorContrastRule().detect(open_pptx(path)))
    assert isinstance(findings, list)


# ---- decorative_flag: rect needs strip geometry ------------------------------


def _pptx_with_rect(tmp_path, cx_in, cy_in, name):
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.enum.shapes import MSO_SHAPE  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    path = tmp_path / name
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, Inches(1), Inches(1), Inches(cx_in), Inches(cy_in)
    )
    pres.save(path)
    return path


def test_normal_empty_rect_not_suggested_decorative(tmp_path):
    path = _pptx_with_rect(tmp_path, 4, 3, "normal_rect.pptx")
    findings = list(DecorativeFlagRule().detect(open_pptx(path)))
    assert findings == []


def test_thin_strip_rect_suggested_decorative(tmp_path):
    path = _pptx_with_rect(tmp_path, 8, 0.2, "strip_rect.pptx")
    findings = list(DecorativeFlagRule().detect(open_pptx(path)))
    assert len(findings) == 1
    assert findings[0].extra["prst"] == "rect"


# ---- list_semantics: manual numbering --------------------------------------


def _docx_with_paragraphs(tmp_path, texts, name):
    from docx import Document  # type: ignore[import-untyped]

    path = tmp_path / name
    docx = Document()
    for t in texts:
        docx.add_paragraph(t)
    docx.save(path)
    return path


def test_manual_numbered_list_detected(tmp_path):
    path = _docx_with_paragraphs(
        tmp_path,
        ["1. first item", "2. second item", "3. third item"],
        "numbered.docx",
    )
    findings = list(ListSemanticsRule().detect(open_docx(path)))
    assert len(findings) == 3


def test_lone_numbered_heading_not_flagged(tmp_path):
    path = _docx_with_paragraphs(
        tmp_path,
        ["1. Introduction", "This chapter introduces the topic at length."],
        "heading.docx",
    )
    findings = list(ListSemanticsRule().detect(open_docx(path)))
    assert findings == []


def test_typed_bullets_still_detected(tmp_path):
    path = _docx_with_paragraphs(
        tmp_path, ["• alpha", "• beta"], "bullets.docx"
    )
    findings = list(ListSemanticsRule().detect(open_docx(path)))
    assert len(findings) == 2


def _docx_with_styled_paragraphs(tmp_path, texts_and_styles, name):
    from docx import Document  # type: ignore[import-untyped]

    path = tmp_path / name
    docx = Document()
    for text, style in texts_and_styles:
        docx.add_paragraph(text, style=style)
    docx.save(path)
    return path


def test_consecutive_numbered_headings_not_flagged(tmp_path):
    """Adjacent numbered section headings must not vouch for each other as a
    fake list — 'promote to a real list' is the wrong fix for headings."""
    path = _docx_with_styled_paragraphs(
        tmp_path,
        [
            ("1. Introduction", "Heading 1"),
            ("2. Methods", "Heading 1"),
            ("3. Results", "Heading 1"),
        ],
        "headings.docx",
    )
    findings = list(ListSemanticsRule().detect(open_docx(path)))
    assert findings == []


def test_numbered_heading_before_typed_bullets_not_flagged(tmp_path):
    """A bullet neighbor must not vouch for a numbered paragraph: a numbered
    heading directly followed by typed bullets is still a heading."""
    path = _docx_with_paragraphs(
        tmp_path,
        ["1. Introduction", "• point one", "• point two"],
        "mixed.docx",
    )
    findings = list(ListSemanticsRule().detect(open_docx(path)))
    # Only the two typed bullets are fake-list items.
    assert len(findings) == 2
    assert all("point" in f.current_value for f in findings)


# ---- reading_order: graphicFrame / grpSp positions ---------------------------


def test_shape_position_reads_graphic_frame_xfrm():
    frame = parse_xml(
        f"""
        <p:graphicFrame {nsdecls("p", "a")}>
          <p:nvGraphicFramePr>
            <p:cNvPr id="7" name="Chart"/>
            <p:cNvGraphicFramePr/>
            <p:nvPr/>
          </p:nvGraphicFramePr>
          <p:xfrm><a:off x="111" y="222"/><a:ext cx="10" cy="10"/></p:xfrm>
          <a:graphic/>
        </p:graphicFrame>
        """
    )
    assert _shape_position(frame) == (222, 111)


def test_shape_position_reads_group_xfrm():
    grp = parse_xml(
        f"""
        <p:grpSp {nsdecls("p", "a")}>
          <p:nvGrpSpPr>
            <p:cNvPr id="8" name="Group"/>
            <p:cNvGrpSpPr/>
            <p:nvPr/>
          </p:nvGrpSpPr>
          <p:grpSpPr>
            <a:xfrm><a:off x="5" y="9"/><a:ext cx="1" cy="1"/></a:xfrm>
          </p:grpSpPr>
        </p:grpSp>
        """
    )
    assert _shape_position(grp) == (9, 5)


def test_shape_position_none_without_explicit_geometry():
    sp = parse_xml(
        f"""
        <p:sp {nsdecls("p", "a")}>
          <p:nvSpPr>
            <p:cNvPr id="9" name="Placeholder"/>
            <p:cNvSpPr/>
            <p:nvPr/>
          </p:nvSpPr>
          <p:spPr/>
        </p:sp>
        """
    )
    assert _shape_position(sp) is None
