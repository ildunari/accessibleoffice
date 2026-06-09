"""Unit test for link_text rule's generic-phrase heuristic."""

from pptx import Presentation  # type: ignore[import-untyped]
from pptx.util import Inches  # type: ignore[import-untyped]

from a11yfix.ooxml.pptx_reader import open_pptx
from a11yfix.rules.link_text import LinkTextRule, _is_generic


def test_click_here_is_generic():
    assert _is_generic("Click here", url=None)


def test_url_only_is_generic():
    assert _is_generic("https://example.com", url=None)


def test_descriptive_text_not_generic():
    assert not _is_generic("Quarterly Earnings Report", url=None)


def test_this_help_center_article_is_generic():
    assert _is_generic("this Help Center article", url="https://help.example.com/article")


def test_ppt_link_in_second_paragraph_uses_second_paragraph_path(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(2))
    shape.text_frame.paragraphs[0].text = "Intro"
    paragraph = shape.text_frame.add_paragraph()
    run = paragraph.add_run()
    run.text = "click here"
    run.hyperlink.address = "https://example.com"
    path = tmp_path / "ppt_link_p2.pptx"
    pres.save(path)

    findings = list(LinkTextRule().detect(open_pptx(path)))

    assert len(findings) == 1
    assert findings[0].officecli_path.endswith("/p[2]/r[1]")
    assert findings[0].officecli_path.startswith("/slide[1]/shape[@id=")
