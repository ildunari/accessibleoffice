"""Unit tests for PPT color contrast detection."""

from __future__ import annotations

from pptx import Presentation  # type: ignore[import-untyped]
from pptx.dml.color import RGBColor  # type: ignore[import-untyped]
from pptx.oxml import parse_xml  # type: ignore[import-untyped]
from pptx.oxml.ns import nsdecls  # type: ignore[import-untyped]
from pptx.util import Inches  # type: ignore[import-untyped]

from a11yfix.ooxml.pptx_reader import open_pptx
from a11yfix.rules.color_contrast import ColorContrastRule


def test_white_text_on_dark_shape_fill_is_not_flagged(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)

    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Readable text"
    run.font.color.rgb = RGBColor(255, 255, 255)

    path = tmp_path / "dark_shape.pptx"
    pres.save(path)

    doc = open_pptx(path)
    findings = list(ColorContrastRule().detect(doc))

    assert findings == []


def test_white_text_on_dark_slide_bg_ref_is_not_flagged(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    bg = parse_xml(
        f'<p:bg {nsdecls("p", "a")}><p:bgRef idx="1001"><a:schemeClr val="dk1"/></p:bgRef></p:bg>'
    )
    slide._element.cSld.insert(0, bg)

    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Readable text"
    run.font.color.rgb = RGBColor(255, 255, 255)

    path = tmp_path / "dark_bg_ref.pptx"
    pres.save(path)

    doc = open_pptx(path)
    findings = list(ColorContrastRule().detect(doc))

    assert findings == []


def test_white_text_on_unknown_background_is_skipped(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Unknown background"
    run.font.color.rgb = RGBColor(255, 255, 255)

    path = tmp_path / "unknown_bg.pptx"
    pres.save(path)

    doc = open_pptx(path)
    findings = list(ColorContrastRule().detect(doc))

    assert findings == []


def test_white_text_on_explicit_white_background_is_flagged(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    bg = parse_xml(
        f'<p:bg {nsdecls("p", "a")}><p:bgPr><a:solidFill><a:srgbClr val="FFFFFF"/>'
        "</a:solidFill></p:bgPr></p:bg>"
    )
    slide._element.cSld.insert(0, bg)

    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    run = shape.text_frame.paragraphs[0].add_run()
    run.text = "Low contrast"
    run.font.color.rgb = RGBColor(255, 255, 255)

    path = tmp_path / "explicit_white_bg.pptx"
    pres.save(path)

    doc = open_pptx(path)
    findings = list(ColorContrastRule().detect(doc))

    assert len(findings) == 1
    assert findings[0].rule_id == "color-contrast"


def test_inherited_paragraph_text_color_is_resolved(tmp_path):
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    shape = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)

    paragraph = shape.text_frame.paragraphs[0]
    p_pr = parse_xml(
        f'<a:pPr {nsdecls("a")}><a:defRPr><a:solidFill><a:srgbClr val="FFFFFF"/>'
        "</a:solidFill></a:defRPr></a:pPr>"
    )
    paragraph._p.insert(0, p_pr)
    run = paragraph.add_run()
    run.text = "Inherited readable text"

    path = tmp_path / "inherited_fg.pptx"
    pres.save(path)

    doc = open_pptx(path)
    findings = list(ColorContrastRule().detect(doc))

    assert findings == []
