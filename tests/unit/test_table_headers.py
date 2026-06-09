"""Unit tests for table_headers rule."""

from pptx.oxml import parse_xml  # type: ignore[import-untyped]
from pptx.oxml.ns import nsdecls  # type: ignore[import-untyped]

from a11yfix.ooxml.docx_reader import open_docx
from a11yfix.ooxml.pptx_reader import open_pptx
from a11yfix.rules.table_headers import TableHeaderRule


def test_table_without_header_detected(docx_no_title):
    doc = open_docx(docx_no_title)
    findings = list(TableHeaderRule().detect(doc))
    assert any(f.rule_id == "table-header-missing" for f in findings)


def test_one_cell_layout_table_not_header_candidate(tmp_path):
    from docx import Document  # type: ignore[import-untyped]

    path = tmp_path / "layout_table.docx"
    docx = Document()
    table = docx.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "layout"
    docx.save(path)

    findings = list(TableHeaderRule().detect(open_docx(path)))

    assert findings == []


def test_plain_ppt_table_missing_header_is_not_auto_fixed(tmp_path):
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    path = tmp_path / "plain_table.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1))
    shape.table.cell(0, 0).text = "plain A"
    shape.table.cell(0, 1).text = "plain B"
    shape._element.graphic.graphicData.tbl.tblPr.set("firstRow", "0")
    pres.save(path)

    doc = open_pptx(path)
    findings = list(TableHeaderRule().detect(doc))

    assert len(findings) == 1
    assert findings[0].extra["visually_header"] is False
    assert TableHeaderRule().fix_deterministic(findings[0], doc) is None


def test_visually_header_ppt_table_is_auto_fixed(tmp_path):
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Inches  # type: ignore[import-untyped]

    path = tmp_path / "header_table.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    shape = slide.shapes.add_table(2, 2, Inches(1), Inches(1), Inches(4), Inches(1))
    shape.table.cell(0, 0).text = "Header A"
    shape.table.cell(0, 1).text = "Header B"
    shape.table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.bold = True
    shape._element.graphic.graphicData.tbl.tblPr.set("firstRow", "0")
    pres.save(path)

    doc = open_pptx(path)
    finding = next(iter(TableHeaderRule().detect(doc)))
    ops = TableHeaderRule().fix_deterministic(finding, doc)

    assert finding.extra["visually_header"] is True
    assert ops is not None
    assert ops[0].props == {"firstRow": "1"}


def test_grouped_ppt_table_header_path_includes_group_scope(tmp_path):
    from pptx import Presentation  # type: ignore[import-untyped]

    path = tmp_path / "grouped_table.pptx"
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    group = parse_xml(
        f"""
        <p:grpSp {nsdecls("p", "a")}>
          <p:nvGrpSpPr>
            <p:cNvPr id="99" name="Grouped objects"/>
            <p:cNvGrpSpPr/>
            <p:nvPr/>
          </p:nvGrpSpPr>
          <p:grpSpPr>
            <a:xfrm>
              <a:off x="0" y="0"/><a:ext cx="1" cy="1"/>
              <a:chOff x="0" y="0"/><a:chExt cx="1" cy="1"/>
            </a:xfrm>
          </p:grpSpPr>
          <p:graphicFrame>
            <p:nvGraphicFramePr>
              <p:cNvPr id="100" name="Grouped table"/>
              <p:cNvGraphicFramePr/>
              <p:nvPr/>
            </p:nvGraphicFramePr>
            <p:xfrm><a:off x="0" y="0"/><a:ext cx="1" cy="1"/></p:xfrm>
            <a:graphic>
              <a:graphicData uri="http://schemas.openxmlformats.org/drawingml/2006/table">
                <a:tbl>
                  <a:tblPr firstRow="0"/>
                  <a:tblGrid><a:gridCol w="1"/><a:gridCol w="1"/></a:tblGrid>
                  <a:tr h="1">
                    <a:tc><a:txBody><a:p><a:r><a:rPr b="true"/><a:t>A</a:t></a:r></a:p></a:txBody></a:tc>
                    <a:tc><a:txBody><a:p><a:r><a:rPr b="true"/><a:t>B</a:t></a:r></a:p></a:txBody></a:tc>
                  </a:tr>
                  <a:tr h="1">
                    <a:tc><a:txBody><a:p><a:r><a:t>1</a:t></a:r></a:p></a:txBody></a:tc>
                    <a:tc><a:txBody><a:p><a:r><a:t>2</a:t></a:r></a:p></a:txBody></a:tc>
                  </a:tr>
                </a:tbl>
              </a:graphicData>
            </a:graphic>
          </p:graphicFrame>
        </p:grpSp>
        """
    )
    slide._element.spTree.append(group)
    pres.save(path)

    doc = open_pptx(path)
    finding = next(iter(TableHeaderRule().detect(doc)))
    ops = TableHeaderRule().fix_deterministic(finding, doc)

    assert finding.officecli_path == "/slide[1]/group[@id=99]/table[@id=100]"
    assert finding.extra["visually_header"] is True
    assert ops is not None
    assert ops[0].path == finding.officecli_path
