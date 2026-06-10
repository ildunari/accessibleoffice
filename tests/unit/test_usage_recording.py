"""Cost recording is the pipeline's job, not the adapter's (one policy, F4)."""
import hashlib
from pathlib import Path

from PIL import Image  # type: ignore[import-untyped]
from pptx import Presentation  # type: ignore[import-untyped]
from pptx.util import Inches  # type: ignore[import-untyped]

from a11yfix.ai.adapter import AltTextResult, CallUsage, LinkTextResult
from a11yfix.cost_meter import _DEFAULT_PRICE_INPUT, CostMeter
from a11yfix.fixers import single_shot
from a11yfix.fixers.single_shot import _record_usage
from a11yfix.ooxml.officecli import BatchResult, ValidationResult


def test_authoritative_cost_recorded(tmp_path: Path, monkeypatch):
    monkeypatch.setenv("A11YFIX_STATE_DIR", str(tmp_path))
    res = AltTextResult(text="t", confidence=0.9, model="some-backend",
                        usage=CallUsage(cost_usd=0.01))
    _record_usage(res)
    assert abs(CostMeter.from_env().total() - 0.01) < 1e-9


def test_token_estimate_fallback(tmp_path: Path, monkeypatch):
    monkeypatch.setenv("A11YFIX_STATE_DIR", str(tmp_path))
    res = AltTextResult(text="t", confidence=0.9, model="unknown-model",
                        usage=CallUsage(input_tokens=1_000_000, output_tokens=0))
    _record_usage(res)
    assert abs(CostMeter.from_env().total() - _DEFAULT_PRICE_INPUT) < 1e-6


def test_cache_tokens_flow_into_estimate(tmp_path: Path, monkeypatch):
    """Cache-read tokens are priced at 10% of the input rate (Finding 2)."""
    monkeypatch.setenv("A11YFIX_STATE_DIR", str(tmp_path))
    res = AltTextResult(
        text="t", confidence=0.9, model="unknown-model",
        usage=CallUsage(cache_read_tokens=1_000_000),
    )
    _record_usage(res)
    assert abs(CostMeter.from_env().total() - _DEFAULT_PRICE_INPUT * 0.10) < 1e-6


def test_no_usage_is_noop(tmp_path: Path, monkeypatch):
    monkeypatch.setenv("A11YFIX_STATE_DIR", str(tmp_path))
    _record_usage(AltTextResult(text="t", confidence=0.9, model="m"))
    assert CostMeter.from_env().total() == 0.0


def test_unwritable_state_dir_degrades_not_raises(tmp_path: Path, monkeypatch):
    ro = tmp_path / "ro"
    ro.mkdir()
    ro.chmod(0o500)  # not writable
    try:
        monkeypatch.setenv("A11YFIX_STATE_DIR", str(ro / "state"))
        meter = CostMeter.from_env()
        # must not raise, must behave as no-op
        meter.record_usd(model="m", usd=0.01)
        assert meter.would_exceed(0.001) is False  # nothing persisted
    finally:
        ro.chmod(0o700)


def test_cache_hit_records_nothing(tmp_path: Path, monkeypatch):
    """A stage-3 cache hit must not touch the cost ledger (Finding 4a):
    the adapter is never called, so no usage exists to record."""
    from a11yfix.ooxml.image_extract import extract_image_for_finding
    from a11yfix.ooxml.pptx_reader import open_pptx
    from a11yfix.rules.alt_text import AltTextRule

    monkeypatch.setenv("A11YFIX_STATE_DIR", str(tmp_path / "state"))
    monkeypatch.setattr(single_shot, "CACHE_DIR", tmp_path / "cache")

    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[6])
    img_path = tmp_path / "image.png"
    Image.new("RGB", (50, 50), color="blue").save(img_path)
    slide.shapes.add_picture(str(img_path), Inches(1), Inches(1), Inches(2), Inches(2))
    deck = tmp_path / "deck.pptx"
    pres.save(deck)

    doc = open_pptx(deck)
    finding = next(iter(AltTextRule().detect(doc)))

    # Pre-populate the cache with exactly the key the pipeline computes.
    ctx = f"Shape: {finding.extra.get('shape_name', finding.extra.get('pic_name', '(unknown)'))}"
    img_bytes, _mime = extract_image_for_finding(doc, finding)
    key = f"alttext|{hashlib.sha256(img_bytes).hexdigest()}|{ctx}"
    single_shot._cache_put(
        key, {"text": "Cached alt text", "confidence": 0.95, "model": "cached"}
    )

    calls = []

    class MeteredAdapter:
        name = "metered"

        def describe_image(self, image_bytes, *, max_chars, context):
            calls.append(context)
            return AltTextResult(
                text="Fresh alt text", confidence=0.95, model="metered",
                usage=CallUsage(cost_usd=1.23),
            )

    class OkClient:
        backup_path = None

        def __init__(self, path, **kwargs):
            self.path = path

        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, tb):
            return None

        def batch(self, ops):
            return BatchResult(success=True, per_op=[{"ok": True} for _ in ops])

        def validate(self):
            return ValidationResult(status="ok")

    monkeypatch.setattr(single_shot, "OfficecliClient", OkClient)

    result = single_shot.apply_single_shot_fixes([finding], doc, MeteredAdapter())

    assert calls == [], "cache hit must not call the adapter"
    assert [fx.after for fx in result.applied] == ["Cached alt text"]
    assert CostMeter.from_env().total() == 0.0


def test_cap_applies_to_any_adapter(tmp_path: Path, monkeypatch):
    """The cost-cap gate is enforced by the pipeline from result.usage, so it
    binds every backend — not just the Claude SDK adapter (F4). The gate runs
    before each adapter call and trips once recorded spend exceeds the cap
    (strictly: meter.total() > cap), so the first 0.30 call against a 0.25
    cap is allowed and the second finding must defer without a model call."""
    from tests.unit.test_fixers_partial_results import FakeDoc, FakeRule, _finding

    monkeypatch.setenv("A11YFIX_STATE_DIR", str(tmp_path / "state"))

    class CostedAdapter:
        name = "costed"

        def __init__(self):
            self.calls = 0

        def suggest_link_text(self, url, surrounding_text):
            self.calls += 1
            return LinkTextResult(
                text="descriptive link", confidence=0.9, model="costed",
                usage=CallUsage(cost_usd=0.30),
            )

    class OkClient:
        backup_path = None

        def __init__(self, path, **kwargs):
            self.path = path

        def __enter__(self):
            return self

        def __exit__(self, exc_type, exc, tb):
            return None

        def batch(self, ops):
            return BatchResult(success=True, per_op=[{"ok": True} for _ in ops])

        def validate(self):
            return ValidationResult(status="ok")

    monkeypatch.setitem(single_shot.REGISTRY, "fake-rule", FakeRule())
    monkeypatch.setattr(single_shot, "OfficecliClient", OkClient)

    doc = FakeDoc(tmp_path / "deck.pptx")
    adapter = CostedAdapter()
    f1, f2 = _finding(1), _finding(2)
    # Distinct surrounding text so f2 can never be satisfied from f1's cache
    # entry — the deferral below must come from the cap gate alone.
    f2.extra["shape_text"] = "tap this"

    result = single_shot.apply_single_shot_fixes(
        [f1, f2], doc, adapter, max_cost_total_usd=0.25
    )

    assert adapter.calls == 1, "cap must stop the second model call"
    assert [fx.finding_id for fx in result.applied] == ["f1"]
    assert abs(CostMeter.from_env().total() - 0.30) < 1e-9
    assert [f.id for f in result.deferred] == ["f2"]
    assert "cost cap" in result.deferred[0].why_human_needed
