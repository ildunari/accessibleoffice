"""Unit tests for slide_title rule."""

from a11yfix.ooxml.pptx_reader import open_pptx
from a11yfix.rules.slide_title import SlideTitleRule


def test_slide_without_title_detected(pptx_no_alt):
    doc = open_pptx(pptx_no_alt)
    findings = list(SlideTitleRule().detect(doc))
    assert findings  # the blank slide has no title
    assert all(f.rule_id == "slide-title-missing" for f in findings)


def test_slide_with_title_not_flagged_for_that_slide(pptx_with_title):
    doc = open_pptx(pptx_with_title)
    findings = list(SlideTitleRule().detect(doc))
    # Slide 2 has the title; slide 1 (blank) does not. Expect ≥1 flagged.
    flagged_indices = {f.extra["slide_index"] for f in findings}
    # The title slide (index 2) should NOT be flagged.
    assert 2 not in flagged_indices


# ---- off-canvas titles (gotcha #10) ------------------------------------------


def _pptx_with_positioned_title(tmp_path, name, *, left, top):
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Emu  # type: ignore[import-untyped]

    path = tmp_path / name
    pres = Presentation()
    slide = pres.slides.add_slide(pres.slide_layouts[0])  # layout with title
    title = slide.shapes.title
    title.text = "Hidden or visible title"
    # Set full geometry (off + ext) the way PowerPoint itself serializes a
    # moved placeholder; position-only writes leave a:ext inherited and the
    # rule conservatively treats that as on-canvas.
    title.left = Emu(left)
    title.top = Emu(top)
    title.width = Emu(6_096_000)  # ~6.67in
    title.height = Emu(914_400)  # 1in
    pres.save(path)
    return path, int(pres.slide_width), int(pres.slide_height)


def test_on_canvas_title_not_flagged(tmp_path):
    path, _w, _h = _pptx_with_positioned_title(
        tmp_path, "visible.pptx", left=914400, top=914400  # 1in x 1in
    )
    findings = list(SlideTitleRule().detect(open_pptx(path)))
    assert findings == []


def test_off_canvas_title_flagged_each_direction(tmp_path):
    cases = {
        "left": (-50_000_000, 914400),
        "top": (914400, -50_000_000),
    }
    for direction, (left, top) in cases.items():
        path, _w, _h = _pptx_with_positioned_title(
            tmp_path, f"off_{direction}.pptx", left=left, top=top
        )
        findings = list(SlideTitleRule().detect(open_pptx(path)))
        assert len(findings) == 1, f"direction={direction}"
        f = findings[0]
        assert f.extra["off_canvas"] is True
        assert f.current_value.startswith("Hidden or visible")


def test_off_canvas_title_beyond_right_and_bottom(tmp_path):
    # Position the box past the slide's far edges.
    path, w, h = _pptx_with_positioned_title(
        tmp_path, "off_right.pptx", left=914400, top=914400
    )
    # Re-open and push beyond the right edge using the known slide width.
    from pptx import Presentation  # type: ignore[import-untyped]
    from pptx.util import Emu  # type: ignore[import-untyped]

    pres = Presentation(str(path))
    title = pres.slides[0].shapes.title
    title.left = Emu(w + 1000)
    pres.save(path)
    findings = list(SlideTitleRule().detect(open_pptx(path)))
    assert len(findings) == 1
    assert findings[0].extra["off_canvas"] is True

    pres = Presentation(str(path))
    title = pres.slides[0].shapes.title
    title.left = Emu(914400)
    title.top = Emu(h + 1000)
    pres.save(path)
    findings = list(SlideTitleRule().detect(open_pptx(path)))
    assert len(findings) == 1


def test_partially_visible_title_not_flagged(tmp_path):
    # Box hangs off the edge but still intersects the slide — visible enough.
    path, _w, _h = _pptx_with_positioned_title(
        tmp_path, "partial.pptx", left=-100_000, top=914400
    )
    findings = list(SlideTitleRule().detect(open_pptx(path)))
    assert findings == []


def test_off_canvas_title_gets_no_single_shot_fix(tmp_path):
    path, _w, _h = _pptx_with_positioned_title(
        tmp_path, "no_fix.pptx", left=-50_000_000, top=914400
    )
    doc = open_pptx(path)
    rule = SlideTitleRule()
    finding = next(iter(rule.detect(doc)))
    assert finding.extra["off_canvas"] is True
    assert rule.fix_single_shot(finding, doc) is None
