"""Pytest fixtures: synthesize .docx and .pptx test files with seeded a11y issues."""

from __future__ import annotations

from pathlib import Path

import pytest
from docx import Document  # type: ignore[import-untyped]
from pptx import Presentation  # type: ignore[import-untyped]
from pptx.util import Inches  # type: ignore[import-untyped]

FIXTURES_DIR = Path(__file__).parent / "fixtures" / "synthetic"
FIXTURES_DIR.mkdir(parents=True, exist_ok=True)


def _write_pptx_basic(path: Path, *, with_alt: bool = False, with_title: bool = True) -> None:
    pres = Presentation()
    blank = pres.slide_layouts[6]  # Blank
    slide = pres.slides.add_slide(blank)
    # Optionally add a title placeholder
    if with_title:
        title_layout = pres.slide_layouts[5]  # Title Only
        title_slide = pres.slides.add_slide(title_layout)
        title_slide.shapes.title.text = "Real Title"
    # Add a picture-shape with no alt text by default
    # python-pptx requires real bytes; use a tiny png
    import io

    from PIL import Image  # type: ignore[import-untyped]

    buf = io.BytesIO()
    Image.new("RGB", (50, 50), color="red").save(buf, format="PNG")
    buf.seek(0)
    pic = slide.shapes.add_picture(buf, Inches(1), Inches(1), Inches(2), Inches(2))
    if with_alt:
        # python-pptx exposes _element to set descr
        nv = pic._element.nvPicPr.cNvPr
        nv.set("descr", "A red square")
    pres.save(str(path))


def _write_docx_basic(path: Path, *, with_title_property: bool = False) -> None:
    doc = Document()
    if with_title_property:
        doc.core_properties.title = "Real Title"
    doc.add_paragraph("Hello world.")
    p = doc.add_paragraph()
    p.add_run("Click here").bold = False
    # add a hyperlink-like generic phrase (not a real hyperlink, but fine for the rule)
    # Add a table without header row
    table = doc.add_table(rows=2, cols=2)
    table.cell(0, 0).text = "A"
    table.cell(0, 1).text = "B"
    table.cell(1, 0).text = "1"
    table.cell(1, 1).text = "2"
    doc.save(str(path))


@pytest.fixture(scope="session")
def pptx_no_alt(tmp_path_factory: pytest.TempPathFactory) -> Path:
    p = tmp_path_factory.mktemp("fixtures") / "pptx_no_alt.pptx"
    _write_pptx_basic(p, with_alt=False, with_title=False)
    return p


@pytest.fixture(scope="session")
def pptx_with_alt(tmp_path_factory: pytest.TempPathFactory) -> Path:
    p = tmp_path_factory.mktemp("fixtures") / "pptx_with_alt.pptx"
    _write_pptx_basic(p, with_alt=True, with_title=True)
    return p


@pytest.fixture(scope="session")
def pptx_with_title(tmp_path_factory: pytest.TempPathFactory) -> Path:
    p = tmp_path_factory.mktemp("fixtures") / "pptx_with_title.pptx"
    _write_pptx_basic(p, with_alt=True, with_title=True)
    return p


@pytest.fixture(scope="session")
def docx_no_title(tmp_path_factory: pytest.TempPathFactory) -> Path:
    p = tmp_path_factory.mktemp("fixtures") / "docx_no_title.docx"
    _write_docx_basic(p, with_title_property=False)
    return p


@pytest.fixture(scope="session")
def docx_with_title(tmp_path_factory: pytest.TempPathFactory) -> Path:
    p = tmp_path_factory.mktemp("fixtures") / "docx_with_title.docx"
    _write_docx_basic(p, with_title_property=True)
    return p
