"""Live end-to-end smoke for one --vlm backend. Costs real money/quota.

Usage:
    uv run python scripts/live_smoke.py --vlm pi [--vlm-model ...]
    uv run python scripts/live_smoke.py --vlm pi --expect-unavailable

Builds a one-slide deck with an image missing alt text, runs stages 1-3,
and asserts the success-gate criteria from the multi-backend plan:
  (a) an AI fix was applied
  (b) the manifest's ai_model matches the backend
  (c) the cost ledger recorded >= 1 call
  (d) a rerun is a cache hit: 0 new calls, byte-identical file

With --expect-unavailable, asserts fail-gate (e) instead: the CLI exits 0,
warns "adapter unavailable" on stderr, and applies no AI fixes.

NOT part of pytest — opt-in only; the controller decides when to spend.
"""

from __future__ import annotations

import argparse
import json
import os
import subprocess
import sys
import tempfile
from pathlib import Path


def build_deck(path: Path) -> None:
    from PIL import Image, ImageDraw
    from pptx import Presentation
    from pptx.util import Inches

    img = path.parent / "chart.png"
    # A describable image, not a flat color: a plain rectangle gets judged
    # DECORATIVE by the model (correctly!) and the fix defers, failing gate (a).
    chart = Image.new("RGB", (320, 200), "white")
    draw = ImageDraw.Draw(chart)
    bars = [("Q1", 60, "#4472c4"), ("Q2", 95, "#ed7d31"), ("Q3", 140, "#70ad47")]
    for i, (label, h, color) in enumerate(bars):
        x = 40 + i * 90
        draw.rectangle([x, 170 - h, x + 60, 170], fill=color, outline="black")
        draw.text((x + 20, 175), label, fill="black")
    draw.line([30, 170, 300, 170], fill="black", width=2)
    draw.text((90, 8), "Quarterly revenue", fill="black")
    chart.save(img)
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    slide.shapes.add_picture(str(img), Inches(1), Inches(1))
    prs.save(path)


def run(deck: Path, state: Path, vlm: str, vlm_model: str | None) -> subprocess.CompletedProcess:
    # A11YFIX_STATE_DIR -> CostMeter writes <state>/cost.json (cost_meter.py).
    # A11YFIX_CACHE must point at a fresh dir: a warm ~/.cache/a11yfix would
    # satisfy stage 3 with 0 model calls and falsely fail gate (c).
    env = {**os.environ, "A11YFIX_STATE_DIR": str(state), "A11YFIX_CACHE": str(state / "cache")}
    # --mode full = stages 1-3 + stage 4; bare "full" would launch an
    # interactive agent session. --dry-run makes stage 4 print its launch
    # plan and return 0 (stage4.launch, cli._launch_stage4 skips the
    # availability check too), so the smoke covers stages 1-3 and terminates
    # without human interaction. The default mode ("auto") would skip stage 3
    # entirely, so it can't be used here.
    # --rules alt-text-missing scopes the run to exactly one finding -> one
    # model call, keeping gates (c)/(d) deterministic (e.g. a missing slide
    # title could otherwise trigger extra, possibly-uncacheable calls).
    cmd = [
        sys.executable, "-m", "a11yfix.cli", str(deck),
        "--mode", "full", "--dry-run",
        "--rules", "alt-text-missing",
        "--vlm", vlm,
        "--output", str(deck.with_suffix(".manifest.json")),
    ]
    if vlm_model:
        cmd += ["--vlm-model", vlm_model]
    return subprocess.run(cmd, env=env, capture_output=True, text=True)


def main() -> int:
    ap = argparse.ArgumentParser(description=__doc__.splitlines()[0])
    ap.add_argument("--vlm", required=True)
    ap.add_argument("--vlm-model", default=None)
    ap.add_argument(
        "--expect-unavailable",
        action="store_true",
        help="Assert the adapter-unavailable path (fail-gate e) instead of the success gates.",
    )
    args = ap.parse_args()

    with tempfile.TemporaryDirectory() as td:
        tmp = Path(td)
        deck = tmp / "smoke.pptx"
        build_deck(deck)

        r1 = run(deck, tmp / "state", args.vlm, args.vlm_model)
        print(r1.stdout[-2000:], r1.stderr[-2000:], sep="\n--- stderr ---\n")
        manifest_path = deck.with_suffix(".manifest.json")
        assert manifest_path.exists(), (
            f"no manifest written (exit {r1.returncode}) — see CLI output above"
        )
        manifest = json.loads(manifest_path.read_text())
        # Stage-3 fixes carry ai_model (= adapter.name); stage-2 fixes never do.
        fixes = [x for x in manifest.get("stage_3_fixes_applied", []) if x.get("ai_model")]

        if args.expect_unavailable:
            # Fail-gate (e): unavailable backend must warn, skip stage 3, exit 0.
            assert r1.returncode == 0, f"GATE (e) FAILED: exit {r1.returncode}, expected 0"
            assert "adapter unavailable" in r1.stderr, (
                "GATE (e) FAILED: stderr lacks 'adapter unavailable' warning"
            )
            assert not fixes, f"GATE (e) FAILED: AI fixes applied despite unavailability: {fixes}"
            print(f"LIVE SMOKE (unavailable path) PASSED for --vlm {args.vlm}")
            return 0

        assert r1.returncode == 0, f"run failed with exit {r1.returncode} — see output above"
        assert fixes, "GATE (a) FAILED: no AI fix applied"
        # ai_model is adapter.name ("pi", "pi:<model>", "codex", ...), so a
        # prefix match works for every backend. Keep the codex escape hatch
        # anyway: its cost-ledger rows use bare model ids, and if ai_model
        # ever follows suit the gate should still mean "codex produced a fix".
        prefix = args.vlm.split(":")[0]
        assert args.vlm == "codex" or any(
            prefix in (x.get("ai_model") or "") for x in fixes
        ), f"GATE (b) FAILED: ai_model {[x['ai_model'] for x in fixes]} does not match backend {args.vlm}"
        ledger = json.loads((tmp / "state" / "cost.json").read_text())
        assert ledger["calls"] >= 1, "GATE (c) FAILED: no cost-ledger calls recorded"

        before = deck.read_bytes()
        r2 = run(deck, tmp / "state", args.vlm, args.vlm_model)
        assert r2.returncode == 0, f"rerun failed with exit {r2.returncode}:\n{r2.stderr[-2000:]}"
        ledger2 = json.loads((tmp / "state" / "cost.json").read_text())
        assert ledger2["calls"] == ledger["calls"], (
            f"GATE (d) FAILED: rerun was not a cache hit "
            f"({ledger2['calls']} calls vs {ledger['calls']})"
        )
        assert deck.read_bytes() == before, "GATE (d) FAILED: rerun changed bytes"

        print(
            f"LIVE SMOKE PASSED for --vlm {args.vlm} "
            f"(cost ${ledger['total_usd']:.4f}, {ledger['calls']} calls)"
        )
    return 0


if __name__ == "__main__":
    raise SystemExit(main())
